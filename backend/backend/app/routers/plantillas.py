from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Query
from fastapi.responses import FileResponse, StreamingResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from pydantic import BaseModel
from typing import Optional, List, Any
import os, uuid, re, io, subprocess, tempfile, shutil
from pathlib import Path
from ..config import UPLOADS_DIR
from ..database import get_db, Plantilla
from ..auth import get_current_user, Usuario
from jose import JWTError, jwt
from ..auth import SECRET_KEY, ALGORITHM

router = APIRouter(prefix="/api/plantillas", tags=["plantillas"])

UPLOAD_DIR = UPLOADS_DIR

FORMATOS_PERMITIDOS = {"pptx", "docx", "pdf"}


def plantilla_dict(p: Plantilla):
    return {
        "id":           p.id,
        "nombre":       p.nombre,
        "evento":       p.evento,
        "fecha_evento": p.fecha_evento,
        "formato":      p.formato,
        "marcadores":   p.marcadores or [],
        "created_at":   p.created_at.isoformat() if p.created_at else None,
    }


def _libreoffice_path() -> Optional[str]:
    """Detecta LibreOffice en Windows, macOS y Linux."""
    candidatos = [
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        # Linux
        "libreoffice",
        "soffice",
    ]
    for c in candidatos:
        if shutil.which(c) or Path(c).exists():
            return c
    return None


def _render_con_libreoffice(ruta: Path, slide_index: int = 0) -> Optional[bytes]:
    """
    Convierte la primera diapositiva/página a PNG usando LibreOffice headless.
    Devuelve bytes del PNG o None si falla.
    """
    lo = _libreoffice_path()
    if not lo:
        return None

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            # LibreOffice convierte a PNG (o PDF intermedio para pptx/docx)
            result = subprocess.run(
                [
                    lo,
                    "--headless",
                    "--norestore",
                    "--nofirststartwizard",
                    "--convert-to", "png",
                    "--outdir", tmpdir,
                    str(ruta),
                ],
                capture_output=True,
                timeout=30,
            )
            # LibreOffice nombra el output igual que el input pero con .png
            # Para multi-página genera nombre-slide1.png, etc.
            tmppath = Path(tmpdir)
            pngs = sorted(tmppath.glob("*.png"))
            if not pngs:
                return None
            # Tomar la primera imagen (primera diapositiva/página)
            return pngs[0].read_bytes()
        except Exception:
            return None


def _render_con_pil_fallback(ruta: Path, formato: str) -> bytes:
    """
    Fallback con PIL: renderiza una vista previa básica pero reconocible
    usando las imágenes incrustadas en el PPTX.
    """
    from PIL import Image as PILImage, ImageDraw
    import io as _io

    try:
        if formato == "pptx":
            from pptx import Presentation
            from pptx.util import Emu

            prs = Presentation(str(ruta))
            if not prs.slides:
                raise ValueError("Sin diapositivas")

            slide = prs.slides[0]
            sw = prs.slide_width
            sh = prs.slide_height
            scale = 1400 / sw
            w = int(sw * scale)
            h = int(sh * scale)

            img = PILImage.new("RGB", (w, h), "white")
            draw = ImageDraw.Draw(img)

            # Ordenar shapes por z-order (primero imágenes, luego texto)
            shapes_sorted = sorted(slide.shapes, key=lambda s: s.shape_type != 13)

            for shape in shapes_sorted:
                # Imágenes incrustadas
                if shape.shape_type == 13:
                    try:
                        sx  = int(shape.left   * scale)
                        sy  = int(shape.top    * scale)
                        sw2 = max(1, int(shape.width  * scale))
                        sh2 = max(1, int(shape.height * scale))
                        pic_bytes = shape.image.blob
                        pic_img = PILImage.open(_io.BytesIO(pic_bytes)).convert("RGBA")
                        pic_img = pic_img.resize((sw2, sh2), PILImage.LANCZOS)
                        img.paste(pic_img, (sx, sy), pic_img)
                    except Exception:
                        pass

                # Fondo de forma con color de relleno
                if hasattr(shape, "fill"):
                    try:
                        fill = shape.fill
                        if fill.type is not None:
                            from pptx.enum.dml import MSO_THEME_COLOR
                            fgcolor = fill.fore_color.rgb
                            sx  = int(shape.left   * scale)
                            sy  = int(shape.top    * scale)
                            sw2 = max(1, int(shape.width  * scale))
                            sh2 = max(1, int(shape.height * scale))
                            draw.rectangle([sx, sy, sx+sw2, sy+sh2],
                                           fill=f"#{fgcolor}", outline=None)
                    except Exception:
                        pass

                # Texto
                if shape.has_text_frame:
                    try:
                        sx = int(shape.left * scale)
                        sy = int(shape.top  * scale)
                        for para in shape.text_frame.paragraphs:
                            texto = para.text.strip()
                            if not texto:
                                continue
                            # Color del texto
                            color = "#222222"
                            try:
                                rgb = para.runs[0].font.color.rgb
                                color = f"#{rgb}"
                            except Exception:
                                pass
                            draw.text((sx + 6, sy + 6), texto[:120], fill=color)
                            sy += 20
                    except Exception:
                        pass

            buf = _io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()

        elif formato == "docx":
            from docx import Document
            doc = Document(str(ruta))
            lines = [p.text for p in doc.paragraphs if p.text.strip()][:30]
            img = PILImage.new("RGB", (800, 1000), "white")
            draw = ImageDraw.Draw(img)
            draw.rectangle([30, 30, 770, 970], outline="#e0e0e0", width=2)
            y = 60
            for line in lines:
                draw.text((50, y), line[:100], fill="#333333")
                y += 28
                if y > 920:
                    break
            buf = _io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()

        else:
            img = PILImage.new("RGB", (600, 800), "#f8f8f8")
            draw = ImageDraw.Draw(img)
            draw.rectangle([20, 20, 580, 780], outline="#cccccc", width=2)
            draw.text((260, 380), "PDF", fill="#aaaaaa")
            buf = _io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()

    except Exception:
        img = PILImage.new("RGB", (400, 560), "#f0f0f0")
        draw = ImageDraw.Draw(img)
        draw.rectangle([10, 10, 390, 550], outline="#cccccc", width=2)
        buf = _io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()


# ── Rutas ────────────────────────────────────────────────────────────────────

@router.get("/")
async def listar_plantillas(
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    result = await db.execute(select(Plantilla).order_by(Plantilla.id.desc()))
    return [plantilla_dict(p) for p in result.scalars().all()]


@router.post("/", status_code=201)
async def subir_plantilla(
    nombre:       str        = Form(...),
    evento:       str        = Form(""),
    fecha_evento: str        = Form(""),
    archivo:      UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    ext = archivo.filename.rsplit(".", 1)[-1].lower()
    if ext not in FORMATOS_PERMITIDOS:
        raise HTTPException(status_code=400, detail=f"Formato no permitido: {ext}")

    filename = f"{uuid.uuid4()}.{ext}"
    ruta = UPLOAD_DIR / filename
    content = await archivo.read()
    ruta.write_bytes(content)

    p = Plantilla(
        nombre=nombre, evento=evento or None,
        fecha_evento=fecha_evento or None,
        formato=ext, ruta_archivo=str(ruta), marcadores=[],
    )
    db.add(p)
    await db.commit()
    await db.refresh(p)
    return plantilla_dict(p)


@router.delete("/{plantilla_id}")
async def eliminar_plantilla(
    plantilla_id: int,
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    result = await db.execute(select(Plantilla).where(Plantilla.id == plantilla_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")
    try:
        Path(p.ruta_archivo).unlink(missing_ok=True)
    except Exception:
        pass
    await db.delete(p)
    await db.commit()
    return {"ok": True}


class MarcadoresUpdate(BaseModel):
    marcadores:   List[Any] = []
    evento:       Optional[str] = None
    fecha_evento: Optional[str] = None


@router.patch("/{plantilla_id}/marcadores")
async def actualizar_marcadores(
    plantilla_id: int,
    data: MarcadoresUpdate,
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    result = await db.execute(select(Plantilla).where(Plantilla.id == plantilla_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")
    p.marcadores = data.marcadores
    if data.evento is not None:
        p.evento = data.evento
    if data.fecha_evento is not None:
        p.fecha_evento = data.fecha_evento
    await db.commit()
    return {"ok": True}


@router.get("/{plantilla_id}/marcadores-detectados")
async def detectar_marcadores(
    plantilla_id: int,
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    result = await db.execute(select(Plantilla).where(Plantilla.id == plantilla_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")

    ruta = Path(p.ruta_archivo)
    if not ruta.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    texto = ""
    try:
        if p.formato == "docx":
            from docx import Document
            doc = Document(str(ruta))
            texto = "\n".join(par.text for par in doc.paragraphs)
        elif p.formato == "pptx":
            from pptx import Presentation
            prs = Presentation(str(ruta))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            texto += " ".join(run.text for run in para.runs) + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al leer plantilla: {e}")

    marcadores = list(set(re.findall(r"§[A-ZÁÉÍÓÚÑÜ_0-9]+", texto)))
    return {"marcadores": sorted(marcadores)}


@router.get("/{plantilla_id}/preview")
async def preview_plantilla(
    plantilla_id: int,
    token: Optional[str] = Query(None),
    db: AsyncSession = Depends(get_db),
):
    # Autenticar via query param (para uso en <img src=...>)
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if not payload.get("sub"):
            raise HTTPException(status_code=401, detail="Not authenticated")
    except JWTError:
        raise HTTPException(status_code=401, detail="Not authenticated")

    result = await db.execute(select(Plantilla).where(Plantilla.id == plantilla_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")

    ruta = Path(p.ruta_archivo)
    if not ruta.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    # 1) Intentar LibreOffice (render fiel)
    img_bytes = _render_con_libreoffice(ruta)

    # 2) Fallback: PIL básico
    if not img_bytes:
        img_bytes = _render_con_pil_fallback(ruta, p.formato)

    return StreamingResponse(io.BytesIO(img_bytes), media_type="image/png")


@router.get("/{plantilla_id}/descargar")
async def descargar_plantilla(
    plantilla_id: int,
    token: Optional[str] = Query(None),
    db: AsyncSession = Depends(get_db),
):
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if not payload.get("sub"):
            raise HTTPException(status_code=401, detail="Not authenticated")
    except JWTError:
        raise HTTPException(status_code=401, detail="Not authenticated")

    result = await db.execute(select(Plantilla).where(Plantilla.id == plantilla_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")

    ruta = Path(p.ruta_archivo)
    if not ruta.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    media_types = {
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pdf":  "application/pdf",
    }
    media_type = media_types.get(p.formato, "application/octet-stream")
    return FileResponse(str(ruta), media_type=media_type, filename=f"{p.nombre}.{p.formato}")