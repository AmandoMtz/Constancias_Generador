from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks, Query, Header
from fastapi.responses import FileResponse, StreamingResponse
import tempfile, shutil
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from pydantic import BaseModel
from typing import List, Optional, Any
import asyncio, zipfile, re, io
from pathlib import Path
from openpyxl import load_workbook
from ..config import LOTES_DIR, UPLOADS_EXCEL_DIR
from ..database import get_db, Plantilla, Persona, Envio
from ..auth import get_current_user, Usuario, SECRET_KEY, ALGORITHM
from jose import JWTError, jwt

router = APIRouter(prefix="/api/constancias", tags=["constancias"])

EXCEL_DIR = UPLOADS_EXCEL_DIR

_estado: dict[int, dict] = {}


class ExcelSource(BaseModel):
    excel_id: str
    sheet_name: str
    nombre_columna: Optional[str] = None


class GenerarRequest(BaseModel):
    plantilla_id: int
    persona_ids: List[int] = []
    excel_source: Optional[ExcelSource] = None
    como_pdf: bool = True
    metodo_envio: str = "zip"
    asunto_email: Optional[str] = None
    cuerpo_email: Optional[str] = None
    datos_extra: Optional[Any] = None


def _excel_path(excel_id: str) -> Optional[Path]:
    posibles = sorted(EXCEL_DIR.glob(f"{excel_id}.*"))
    return posibles[0] if posibles else None


def _clean_str(val: Any) -> str:
    return "" if val is None else str(val)


def _normalize_match(val: Any) -> str:
    texto = _clean_str(val).strip().lower()
    return " ".join(texto.split())


def _row_value(source: Any, key: str, datos_extra: dict | None = None) -> str:
    if not key:
        return ""
    if isinstance(source, dict):
        if key in source:
            return _clean_str(source.get(key))
    else:
        if hasattr(source, key):
            return _clean_str(getattr(source, key, ""))
    if datos_extra:
        return _clean_str(datos_extra.get(key, ""))
    return ""


def _nombre_corregido(nombre: str) -> str:
    nombre = _clean_str(nombre).strip()
    if not nombre:
        return ""
    patrones = [
        r"^dr\.\s*", r"^dra\.\s*", r"^m\.c\.\s*", r"^m\.a\.\s*",
        r"^ing\.\s*", r"^lic\.\s*", r"^mtro\.\s*", r"^mtra\.\s*",
        r"^arq\.\s*",
    ]
    for pat in patrones:
        nombre = re.sub(pat, "", nombre, flags=re.IGNORECASE)
    return " ".join(nombre.split())


def _tratamiento_dr(base: str) -> str:
    base = _clean_str(base).lower()
    if any(x in base for x in [" dra", "dra.", "doctora", "mtra", "maestra", "arquitecta", "licenciada", "ingeniera"]):
        return "Dra."
    return "Dr."


def _tratamiento_c(base: str) -> str:
    base = _clean_str(base).lower()
    if any(x in base for x in [" dra", "dra.", "doctora", "mtra", "maestra", "arquitecta", "licenciada", "ingeniera"]):
        return "A la C."
    return "Al C."


def _extract_token(
    token: Optional[str] = None,
    authorization: Optional[str] = None,
) -> str:
    if token:
        return token

    if authorization:
        auth = authorization.strip()
        if auth.lower().startswith("bearer "):
            return auth[7:].strip()
        return auth

    raise HTTPException(status_code=401, detail="Not authenticated")


def _validate_token(
    token: Optional[str] = None,
    authorization: Optional[str] = None,
) -> dict:
    raw_token = _extract_token(token=token, authorization=authorization)
    try:
        payload = jwt.decode(raw_token, SECRET_KEY, algorithms=[ALGORITHM])
        if not payload.get("sub"):
            raise HTTPException(status_code=401, detail="Not authenticated")
        return payload
    except JWTError:
        raise HTTPException(status_code=401, detail="Not authenticated")


def reemplazar_marcadores(texto: str, source, marcadores: list, datos_extra: dict | None) -> str:
    resultado = texto
    datos_extra = datos_extra or {}

    marcadores_ordenados = sorted(
        marcadores,
        key=lambda m: len((m.get("marcador") or "")),
        reverse=True,
    )

    for m in marcadores_ordenados:
        marcador = m.get("marcador", "")
        tipo = m.get("tipo", "columna")
        columna = m.get("columna", "")
        valor = m.get("valor", "")
        cols_join = m.get("columnas", "")
        base = _row_value(source, columna, datos_extra)

        if tipo == "texto":
            reemplazo = _clean_str(valor)
        elif tipo == "nombre":
            reemplazo = _nombre_corregido(base)
        elif tipo == "trat_dr":
            reemplazo = _tratamiento_dr(base)
        elif tipo == "trat_c":
            reemplazo = _tratamiento_c(base)
        elif tipo == "columnas_join":
            partes = [col.strip() for col in _clean_str(cols_join).split(",") if col.strip()]
            vals = [_row_value(source, col, datos_extra).strip() for col in partes]
            reemplazo = " ".join(v for v in vals if v)
        else:
            reemplazo = base

        resultado = resultado.replace(marcador, _clean_str(reemplazo))

    for key in ["evento", "fecha"]:
        resultado = resultado.replace(f"§{key.upper()}", _clean_str(datos_extra.get(key, "")))

    return resultado


def procesar_docx(ruta, source, marcadores, datos_extra):
    from docx import Document

    doc = Document(str(ruta))

    def replace_in_para(para):
        texto_original = para.text
        nuevo_texto = reemplazar_marcadores(texto_original, source, marcadores, datos_extra or {})
        if nuevo_texto == texto_original:
            return
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = nuevo_texto
        else:
            para.add_run(nuevo_texto)

    for para in doc.paragraphs:
        replace_in_para(para)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    replace_in_para(para)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def procesar_pptx(ruta, source, marcadores, datos_extra):
    from pptx import Presentation

    prs = Presentation(str(ruta))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                cambio_run = False
                for run in para.runs:
                    original = run.text
                    nuevo = reemplazar_marcadores(original, source, marcadores, datos_extra or {})
                    if nuevo != original:
                        run.text = nuevo
                        cambio_run = True

                if cambio_run:
                    continue

                texto_full = "".join(run.text for run in para.runs)
                nuevo_texto = reemplazar_marcadores(texto_full, source, marcadores, datos_extra or {})

                if nuevo_texto != texto_full and para.runs:
                    para.runs[0].text = nuevo_texto
                    for run in para.runs[1:]:
                        run.text = ""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _convertir_a_pdf_bytes(input_path: Path) -> bytes:
    import subprocess
    import traceback

    with tempfile.TemporaryDirectory(prefix="const_pdf_") as tmpdir:
        outdir = Path(tmpdir)

        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", str(outdir),
                    str(input_path),
                ],
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice error: {result.stderr}")

        except FileNotFoundError:
            raise RuntimeError("LibreOffice no está instalado en el servidor")
        except Exception as e:
            raise RuntimeError(f"Error al convertir con LibreOffice: {e}\n{traceback.format_exc()}")

        pdf_path = outdir / f"{input_path.stem}.pdf"
        if not pdf_path.exists():
            raise RuntimeError(f"No se generó el PDF esperado: {pdf_path}")

        return pdf_path.read_bytes()


def _generar_pdf_desde_plantilla(ruta: Path, formato: str, source, marcadores, datos_extra) -> bytes:
    with tempfile.TemporaryDirectory(prefix="const_tmp_") as tmpdir:
        tmpdir_path = Path(tmpdir)

        if formato == "docx":
            procesado = procesar_docx(ruta, source, marcadores, datos_extra)
            entrada = tmpdir_path / "documento.docx"
        elif formato == "pptx":
            procesado = procesar_pptx(ruta, source, marcadores, datos_extra)
            entrada = tmpdir_path / "presentacion.pptx"
        else:
            raise RuntimeError(f"Formato no soportado para PDF: {formato}")

        entrada.write_bytes(procesado)
        return _convertir_a_pdf_bytes(entrada)


def _generar_thumbnail_desde_archivo(archivo: Path) -> bytes:
    from PIL import Image as PILImage, ImageDraw
    ext = archivo.suffix.lower().lstrip(".")

    try:
        if ext == "pdf":
            import fitz

            doc = fitz.open(str(archivo))
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
            doc.close()

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(str(archivo))
            if not prs.slides:
                raise ValueError("Sin diapositivas")

            slide = prs.slides[0]
            sw = prs.slide_width
            sh = prs.slide_height
            scale = 900 / sw
            w = int(sw * scale)
            h = int(sh * scale)
            img = PILImage.new("RGB", (w, h), "white")
            draw = ImageDraw.Draw(img)

            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        sx = int(shape.left * scale)
                        sy = int(shape.top * scale)
                        sw2 = max(1, int(shape.width * scale))
                        sh2 = max(1, int(shape.height * scale))
                        pic_img = PILImage.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                        pic_img = pic_img.resize((sw2, sh2), PILImage.LANCZOS)
                        img.paste(pic_img, (sx, sy), pic_img)
                    except Exception:
                        pass

                if shape.has_text_frame:
                    sx = int(shape.left * scale)
                    sy = int(shape.top * scale)
                    sw2 = max(1, int(shape.width * scale))
                    sh2 = max(1, int(shape.height * scale))
                    draw.rectangle([sx, sy, sx + sw2, sy + sh2], outline="#e8e8e8", width=1)
                    texto = shape.text_frame.text[:300]
                    if texto.strip():
                        try:
                            draw.text((sx + 6, sy + 6), texto, fill="#222222")
                        except Exception:
                            pass

        elif ext == "docx":
            from docx import Document
            doc = Document(str(archivo))
            lines = [p.text for p in doc.paragraphs if p.text.strip()][:25]
            img = PILImage.new("RGB", (800, 1000), "white")
            draw = ImageDraw.Draw(img)
            draw.rectangle([30, 30, 770, 970], outline="#e0e0e0", width=2)
            y = 60
            for line in lines:
                draw.text((50, y), line[:100], fill="#333333")
                y += 28
                if y > 920:
                    break

        else:
            img = PILImage.new("RGB", (600, 800), "#f8f8f8")
            draw = ImageDraw.Draw(img)
            draw.rectangle([20, 20, 580, 780], outline="#cccccc", width=2)
            draw.text((260, 370), "PDF", fill="#aaaaaa")

    except Exception:
        img = PILImage.new("RGB", (400, 560), "#f0f0f0")
        draw = ImageDraw.Draw(img)
        draw.rectangle([10, 10, 390, 550], outline="#cccccc", width=2)

    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()

def _iter_excel_rows(excel_source: ExcelSource):
    ruta = _excel_path(excel_source.excel_id)
    if not ruta:
        raise HTTPException(status_code=404, detail="Excel no encontrado")

    wb = load_workbook(ruta, read_only=True, data_only=True)
    if excel_source.sheet_name not in wb.sheetnames:
        wb.close()
        raise HTTPException(status_code=404, detail="Hoja de Excel no encontrada")

    ws = wb[excel_source.sheet_name]
    headers = []
    for idx, cell in enumerate(ws[1], start=1):
        val = cell.value
        headers.append(str(val).strip() if val not in (None, "") else f"col_{idx}")

    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        item = {headers[i]: row[i] if i < len(row) else None for i in range(len(headers))}
        if any(v not in (None, "") for v in item.values()):
            rows.append(item)

    wb.close()
    return rows


def _filtrar_excel_rows_por_personas(rows: list[dict], personas: list[Persona], nombre_columna: str | None) -> list[dict]:
    if not personas:
        return rows

    nombres = {_normalize_match(p.nombre) for p in personas if getattr(p, "nombre", None)}
    correos = {_normalize_match(p.email) for p in personas if getattr(p, "email", None)}
    matriculas = {_normalize_match(p.matricula) for p in personas if getattr(p, "matricula", None)}

    filtrados = []
    for row in rows:
        candidatos = set()
        if nombre_columna:
            candidatos.add(_normalize_match(row.get(nombre_columna)))
        candidatos.add(_normalize_match(row.get("nombre")))
        candidatos.add(_normalize_match(row.get("email")))
        candidatos.add(_normalize_match(row.get("matricula")))
        candidatos.discard("")

        if (nombres and candidatos & nombres) or (correos and candidatos & correos) or (matriculas and candidatos & matriculas):
            filtrados.append(row)

    return filtrados


async def generar_lote(envio_id, plantilla, personas, metodo, marcadores, datos_extra, asunto, cuerpo, excel_source, persona_ids=None):
    _estado[envio_id]["estado"] = "procesando"

    lote_dir = LOTES_DIR / str(envio_id)
    lote_dir.mkdir(parents=True, exist_ok=True)

    ruta = Path(plantilla.ruta_archivo)
    formato = plantilla.formato
    enviados = 0
    errores = 0

    fuentes = personas
    nombre_columna = excel_source.get("nombre_columna") if excel_source else None
    if excel_source:
        fuentes = _iter_excel_rows(ExcelSource(**excel_source))
        if persona_ids:
            fuentes = _filtrar_excel_rows_por_personas(fuentes, personas, nombre_columna)

    for idx, source in enumerate(fuentes, start=1):
        try:
            if formato in {"docx", "pptx"}:
                if _estado[envio_id].get("como_pdf", True):
                    data = _generar_pdf_desde_plantilla(ruta, formato, source, marcadores, datos_extra)
                    ext_out = "pdf"
                elif formato == "docx":
                    data = procesar_docx(ruta, source, marcadores, datos_extra)
                    ext_out = "docx"
                else:
                    data = procesar_pptx(ruta, source, marcadores, datos_extra)
                    ext_out = "pptx"
            else:
                data = ruta.read_bytes()
                ext_out = "pdf"

            if isinstance(source, dict):
                nombre_base = _row_value(source, nombre_columna or "", datos_extra) if nombre_columna else ""
                if not nombre_base:
                    nombre_base = _row_value(source, "nombre", datos_extra)
                if not nombre_base:
                    nombre_base = _row_value(source, "nombre_evaluador", datos_extra)
                if not nombre_base:
                    nombre_base = f"registro_{idx}"
            else:
                nombre_base = getattr(source, "nombre", None) or f"persona_{getattr(source, 'id', idx)}"

            nombre_archivo = re.sub(r"[^\w\-_\. ]", "_", _clean_str(nombre_base)).strip()
            if not nombre_archivo:
                nombre_archivo = f"registro_{idx}"

            base_filename = f"Constancia_{nombre_archivo}"
            final_path = lote_dir / f"{base_filename}.{ext_out}"

            contador = 2
            while final_path.exists():
                final_path = lote_dir / f"{base_filename}_{contador}.{ext_out}"
                contador += 1

            final_path.write_bytes(data)
            enviados += 1

        except Exception as e:
            errores += 1
            print(f"ERROR generando constancia {idx}: {e}")

        _estado[envio_id]["enviados"] = enviados
        _estado[envio_id]["errores"] = errores
        await asyncio.sleep(0.03)

    zip_path = None
    if metodo == "zip":
        zip_path = LOTES_DIR / f"lote_{envio_id}.zip"
        with zipfile.ZipFile(str(zip_path), "w", zipfile.ZIP_DEFLATED) as zf:
            for archivo in lote_dir.iterdir():
                zf.write(archivo, archivo.name)

    _estado[envio_id]["estado"] = "completado" if errores == 0 else "error"
    _estado[envio_id]["zip_listo"] = zip_path is not None and zip_path.exists()
    _estado[envio_id]["zip_path"] = str(zip_path) if zip_path else None

    from ..database import AsyncSessionLocal
    async with AsyncSessionLocal() as session:
        result = await session.execute(select(Envio).where(Envio.id == envio_id))
        envio = result.scalar_one_or_none()
        if envio:
            envio.estado = _estado[envio_id]["estado"]
            envio.enviados = enviados
            envio.errores = errores
            envio.zip_listo = _estado[envio_id]["zip_listo"]
            envio.ruta_zip = _estado[envio_id].get("zip_path")
            await session.commit()

@router.post("/generar")
async def generar_constancias(
    req: GenerarRequest,
    background_tasks: BackgroundTasks,
    db: AsyncSession = Depends(get_db),
    _: Usuario = Depends(get_current_user),
):
    result = await db.execute(select(Plantilla).where(Plantilla.id == req.plantilla_id))
    plantilla = result.scalar_one_or_none()
    if not plantilla:
        raise HTTPException(status_code=404, detail="Plantilla no encontrada")

    personas = []
    total = 0
    excel_source_payload = None

    if req.persona_ids:
        result = await db.execute(select(Persona).where(Persona.id.in_(req.persona_ids)))
        personas = result.scalars().all()
        if not personas:
            raise HTTPException(status_code=400, detail="No se encontraron destinatarios seleccionados")

    if req.excel_source:
        rows = _iter_excel_rows(req.excel_source)
        if not rows:
            raise HTTPException(status_code=400, detail="La hoja de Excel no contiene registros")

        excel_source_payload = req.excel_source.model_dump()
        if personas:
            rows = _filtrar_excel_rows_por_personas(rows, personas, req.excel_source.nombre_columna)
            if not rows:
                raise HTTPException(status_code=400, detail="Ningún alumno seleccionado coincide con los registros del Excel")
        total = len(rows)
    else:
        total = len(personas)
        if not personas:
            raise HTTPException(status_code=400, detail="No se encontraron destinatarios")

    envio = Envio(
        plantilla_id=req.plantilla_id,
        metodo=req.metodo_envio,
        estado="pendiente",
        total=total,
        enviados=0,
        errores=0,
    )
    db.add(envio)
    await db.commit()
    await db.refresh(envio)

    _estado[envio.id] = {
        "estado": "pendiente",
        "total": total,
        "enviados": 0,
        "errores": 0,
        "zip_listo": False,
        "zip_path": None,
        "como_pdf": req.como_pdf,
    }

    background_tasks.add_task(
        generar_lote,
        envio.id,
        plantilla,
        list(personas),
        req.metodo_envio,
        plantilla.marcadores or [],
        req.datos_extra or {},
        req.asunto_email,
        req.cuerpo_email,
        excel_source_payload,
        req.persona_ids,
    )
    return {"envio_id": envio.id, "total": total}


@router.get("/estado/{envio_id}")
async def estado_envio(envio_id: int, _: Usuario = Depends(get_current_user)):
    if envio_id in _estado:
        return _estado[envio_id]

    from ..database import AsyncSessionLocal
    async with AsyncSessionLocal() as session:
        result = await session.execute(select(Envio).where(Envio.id == envio_id))
        envio = result.scalar_one_or_none()
        if not envio:
            raise HTTPException(status_code=404, detail="Envío no encontrado")

        return {
            "estado": envio.estado,
            "total": envio.total,
            "enviados": envio.enviados,
            "errores": envio.errores,
            "zip_listo": envio.zip_listo,
        }


@router.get("/lote/{envio_id}/previews")
async def previews_lote(
    envio_id: int,
    token: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    _validate_token(token=token, authorization=authorization)

    lote_dir = LOTES_DIR / str(envio_id)
    if not lote_dir.exists():
        raise HTTPException(status_code=404, detail="Lote no encontrado")

    archivos = sorted(lote_dir.iterdir())
    if not archivos:
        return []

    return [
        {
            "nombre_archivo": archivo.name,
            "url_preview": f"/api/constancias/lote/{envio_id}/preview/{archivo.name}",
            "url_archivo": f"/api/constancias/lote/{envio_id}/archivo/{archivo.name}",
        }
        for archivo in archivos
        if archivo.is_file()
    ]


@router.get("/lote/{envio_id}/preview/{nombre_archivo}")
async def preview_constancia(
    envio_id: int,
    nombre_archivo: str,
    token: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    _validate_token(token=token, authorization=authorization)

    archivo = LOTES_DIR / str(envio_id) / nombre_archivo
    if not archivo.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    try:
        img_bytes = _generar_thumbnail_desde_archivo(archivo)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generando thumbnail: {e}")

    return StreamingResponse(io.BytesIO(img_bytes), media_type="image/png")


@router.get("/lote/{envio_id}/zip")
async def descargar_zip(
    envio_id: int,
    token: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    _validate_token(token=token, authorization=authorization)

    zip_path = None
    if envio_id in _estado:
        zip_path = _estado[envio_id].get("zip_path")

    if not zip_path:
        from ..database import AsyncSessionLocal
        async with AsyncSessionLocal() as session:
            result = await session.execute(select(Envio).where(Envio.id == envio_id))
            envio = result.scalar_one_or_none()
            if envio:
                zip_path = envio.ruta_zip

    if not zip_path or not Path(zip_path).exists():
        raise HTTPException(status_code=404, detail="ZIP no disponible")

    return FileResponse(
        zip_path,
        media_type="application/zip",
        filename=f"constancias_lote_{envio_id}.zip",
    )


@router.get("/lote/{envio_id}/archivo/{nombre_archivo}")
async def descargar_archivo_generado(
    envio_id: int,
    nombre_archivo: str,
    token: Optional[str] = Query(None),
    authorization: Optional[str] = Header(None),
):
    _validate_token(token=token, authorization=authorization)

    archivo = LOTES_DIR / str(envio_id) / nombre_archivo
    if not archivo.exists():
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    media = {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }.get(archivo.suffix.lower(), "application/octet-stream")

    return FileResponse(
    str(archivo),
    media_type=media,
    filename=archivo.name,
    content_disposition_type="inline",
)
